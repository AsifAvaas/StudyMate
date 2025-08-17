import os
import asyncio
import tempfile
from sqlalchemy.orm import Session
from fastapi import UploadFile, File, HTTPException, APIRouter,Depends,BackgroundTasks
from dotenv import load_dotenv
from PIL import Image
import pytesseract
from utils.Cloudinary_config import cloudinary
from cloudinary.uploader import upload
import pdfplumber
from pdf2image import convert_from_path
from pptx import Presentation
from docx import Document
from models.Resources import resources
from database import get_db
from utils.Process_Chunk import process_chunks

load_dotenv()

router = APIRouter(prefix="/cloud", tags=["Resources"])


ALLOWED_TYPES = [
    "application/pdf",  # PDF
    "application/msword",  # DOC
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # DOCX
    "application/vnd.ms-powerpoint",  # PPT
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # PPTX
]



def extract_text(file_path, file_type):
    text = ""
    
    if file_type == "application/pdf":
        # Extract text with pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        # OCR fallback if no text found
        if not text.strip():
            images = convert_from_path(file_path)
            for img in images:
                text += pytesseract.image_to_string(img)
    
    elif file_type in ["application/msword",
                       "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
        
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        # Optional: extract images for OCR
        for rel in doc.part._rels:
            rel = doc.part._rels[rel]
            if "image" in rel.target_ref:
                img_data = rel.target_part.blob
                img = Image.open(io.BytesIO(img_data))
                text += "\n" + pytesseract.image_to_string(img)
    
    elif file_type in ["application/vnd.ms-powerpoint",
                       "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text += para.text + "\n"
                # Extract images for OCR
                if hasattr(shape, "image"):
                    img = shape.image
                    pil_img = Image.open(io.BytesIO(img.blob))
                    text += "\n" + pytesseract.image_to_string(pil_img)
    
    return text.strip()




# def process_chunks(db: Session, file_id: str, user_id: str, text: str):
#     Chunk_data = chunk_text(text, max_words=200)
#     for ch in Chunk_data:
#         embedding = get_embedding(ch)
#         new_chunk = chunks(
#             file_id=file_id,
#             user_id=user_id,
#             chunk_text=ch,
#             embedding=embedding
#         )
#         db.add(new_chunk)
#     db.commit()
#     auto_segment_topics(file_id,db)



@router.post('/upload/{user_id}')
async def upload_to_cloud(
    user_id: str,
    file: UploadFile = File(...),
    db: Session= Depends(get_db),
    background_tasks: BackgroundTasks=None
    ):
    if file.content_type not in ALLOWED_TYPES:
        raise HTTPException(status_code=400, detail="Invalid file type")
    
    # Save temp file
    suffix = os.path.splitext(file.filename)[1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    # Extract text automatically
    extracted_text = extract_text(tmp_path, file.content_type)

    # Upload to Cloudinary
    file_name, file_ext = os.path.splitext(file.filename)
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(
        None,
        lambda: upload(
            tmp_path,
            folder="Studymate/resources",
            resource_type="raw",
            public_id=file_name,
            overwrite=True,
            format=file_ext.lstrip(".")
        )
    )

    os.remove(tmp_path)
    new_resource= resources(
        resource_url=result["secure_url"],
        user_id=user_id,
        resource_text=extracted_text
    )
    db.add(new_resource)
    db.commit()
    db.refresh(new_resource)
    # background_tasks.add_task(process_chunks,db,new_resource.id,user_id,extracted_text)
    background_tasks.add_task(process_chunks, new_resource.id, user_id, extracted_text)


    return {"url": result["secure_url"]}
